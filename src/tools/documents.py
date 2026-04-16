"""
documents.py — Generazione documenti: PDF, DOCX, XLSX, PPTX

Tutti i file vengono salvati in data/output/ e il percorso viene restituito
all'agent core, che lo passerà al bot handler per l'invio su Telegram.
"""

from datetime import datetime
from pathlib import Path

from loguru import logger

OUTPUT_DIR = Path("data/output")


def _ensure_output_dir() -> None:
    OUTPUT_DIR.mkdir(parents=True, exist_ok=True)


def _output_path(filename: str) -> Path:
    _ensure_output_dir()
    return OUTPUT_DIR / filename


def _timestamped_name(prefix: str, ext: str, filename: str | None) -> str:
    if filename and filename.endswith(f".{ext}"):
        return filename
    if filename:
        return f"{filename}.{ext}"
    return f"{prefix}_{datetime.now().strftime('%Y%m%d_%H%M%S')}.{ext}"


# ─── Utility ──────────────────────────────────────────────────────────────────

def _clean_text(text: str) -> str:
    """
    Sostituisce i caratteri Unicode "fantasiosi" (smart quotes, em-dashes) 
    con equivalenti ASCII per evitare crash nei PDF con font standard.
    """
    if not text:
        return ""
    replacements = {
        '\u2013': '-', # en-dash
        '\u2014': '-', # em-dash
        '\u2018': "'", # left single quote
        '\u2019': "'", # right single quote
        '\u201c': '"', # left double quote
        '\u201d': '"', # right double quote
        '\u2026': '...', # ellipsis
        '\u00a0': ' ',   # non-breaking space
    }
    for char, replacement in replacements.items():
        text = text.replace(char, replacement)
    
    # Fallback finale: rimuove qualsiasi carattere non latin-1 per evitare crash
    return text.encode('latin-1', 'replace').decode('latin-1')


# ─── PDF ──────────────────────────────────────────────────────────────────────

def generate_pdf(title: str, content: str, filename: str | None = None) -> dict:
    """
    Genera un documento PDF con titolo e contenuto testuale.
    """
    try:
        from fpdf import FPDF

        fname = _timestamped_name("documento", "pdf", filename)
        out = _output_path(fname)

        # Pulizia testo per compatibilità font standard (Helvetica)
        clean_title = _clean_text(title)
        clean_content = _clean_text(content)

        pdf = FPDF()
        pdf.set_margins(20, 20, 20)
        pdf.add_page()

        # Titolo
        pdf.set_font("Helvetica", "B", 18)
        pdf.multi_cell(0, 12, clean_title, align="C")
        pdf.ln(6)

        # Linea separatrice
        pdf.set_draw_color(100, 100, 100)
        pdf.line(20, pdf.get_y(), 190, pdf.get_y())
        pdf.ln(6)

        # Corpo
        pdf.set_font("Helvetica", size=11)
        pdf.multi_cell(0, 7, clean_content)

        # Footer
        pdf.set_y(-15)
        pdf.set_font("Helvetica", "I", 8)
        pdf.set_text_color(150, 150, 150)
        pdf.cell(0, 10, f"Generato da Demuclaw - {datetime.now().strftime('%d/%m/%Y %H:%M')}", align="C")

        pdf.output(str(out))
        logger.info(f"📄 PDF generato: {out}")
        return {"success": True, "path": str(out), "filename": fname}

    except Exception as e:
        logger.error(f"Errore generazione PDF: {e}")
        return {"error": str(e)}


# ─── DOCX ─────────────────────────────────────────────────────────────────────

def generate_docx(title: str, content: str, filename: str | None = None) -> dict:
    """
    Genera un documento Word (.docx) con titolo e contenuto.

    Args:
        title:   Titolo (Heading 1).
        content: Corpo (i doppi newline diventano paragrafi separati).
    """
    try:
        from docx import Document
        from docx.shared import Pt, RGBColor

        fname = _timestamped_name("documento", "docx", filename)
        out = _output_path(fname)

        doc = Document()

        # Pulizia testo
        clean_title = _clean_text(title)
        clean_content = _clean_text(content)

        # Titolo
        heading = doc.add_heading(clean_title, level=0)
        heading.runs[0].font.color.rgb = RGBColor(0x1F, 0x4E, 0x79)

        # Paragrafi (divisi per doppio newline)
        paragraphs = [p.strip() for p in clean_content.split("\n\n") if p.strip()]
        for para_text in paragraphs:
            p = doc.add_paragraph(para_text)
            p.style.font.size = Pt(11)

        # Footer
        doc.add_paragraph(
            f"\nGenerato da Demuclaw — {datetime.now().strftime('%d/%m/%Y %H:%M')}"
        ).runs[0].font.color.rgb = RGBColor(0x80, 0x80, 0x80)

        doc.save(str(out))
        logger.info(f"📝 DOCX generato: {out}")
        return {"success": True, "path": str(out), "filename": fname}

    except Exception as e:
        logger.error(f"Errore generazione DOCX: {e}")
        return {"error": str(e)}


# ─── XLSX ─────────────────────────────────────────────────────────────────────

def generate_xlsx(
    title: str,
    data: list,
    headers: list | None = None,
    filename: str | None = None,
) -> dict:
    """
    Genera un foglio Excel (.xlsx).

    Args:
        title:   Titolo del foglio (usato come nome del tab).
        data:    Lista di righe; ogni riga è una lista di valori.
        headers: Intestazioni delle colonne (opzionali, con sfondo blu).
        filename: Nome del file output.

    Esempio:
        generate_xlsx(
            title="Spese",
            headers=["Data", "Descrizione", "Importo"],
            data=[["01/01/2025", "Affitto", 800], ["05/01/2025", "Spesa", 120]]
        )
    """
    try:
        from openpyxl import Workbook
        from openpyxl.styles import Alignment, Font, PatternFill

        fname = _timestamped_name("foglio", "xlsx", filename)
        out = _output_path(fname)

        wb = Workbook()
        ws = wb.active
        clean_title_sheet = _clean_text(title)[:31]
        ws.title = clean_title_sheet # Excel: max 31 caratteri per nome tab

        # Stile header
        header_fill = PatternFill(start_color="1F4E79", end_color="1F4E79", fill_type="solid")
        header_font = Font(color="FFFFFF", bold=True, size=11)
        header_align = Alignment(horizontal="center", vertical="center")

        start_row = 1
        if headers:
            for col_idx, header in enumerate(headers, 1):
                cell = ws.cell(row=1, column=col_idx, value=header)
                cell.fill = header_fill
                cell.font = header_font
                cell.alignment = header_align
                ws.column_dimensions[cell.column_letter].width = max(15, len(str(header)) + 4)
            start_row = 2

        # Dati
        for row_idx, row_data in enumerate(data, start_row):
            for col_idx, value in enumerate(row_data, 1):
                clean_value = _clean_text(str(value)) if isinstance(value, str) else value
                ws.cell(row=row_idx, column=col_idx, value=clean_value)

        wb.save(str(out))
        logger.info(f"📊 XLSX generato: {out}")
        return {"success": True, "path": str(out), "filename": fname}

    except Exception as e:
        logger.error(f"Errore generazione XLSX: {e}")
        return {"error": str(e)}


# ─── PPTX ─────────────────────────────────────────────────────────────────────

def generate_pptx(
    title: str,
    slides: list,
    filename: str | None = None,
) -> dict:
    """
    Genera una presentazione PowerPoint (.pptx).

    Args:
        title:  Titolo della presentazione (slide di apertura).
        slides: Lista di dict {"title": str, "content": str}.
        filename: Nome del file output.

    Esempio:
        generate_pptx(
            title="Presentazione Progetto",
            slides=[
                {"title": "Introduzione", "content": "Descrizione..."},
                {"title": "Roadmap", "content": "Fase 1: ..."},
            ]
        )
    """
    try:
        from pptx import Presentation
        from pptx.dml.color import RGBColor
        from pptx.util import Pt

        fname = _timestamped_name("presentazione", "pptx", filename)
        out = _output_path(fname)

        prs = Presentation()

        # ── Slide titolo ──────────────────────────────────────
        title_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(title_layout)
        slide.shapes.title.text = _clean_text(title)
        if len(slide.placeholders) > 1:
            slide.placeholders[1].text = (
                f"Generato da Demuclaw - {datetime.now().strftime('%d/%m/%Y')}"
            )

        # ── Slide contenuto ───────────────────────────────────
        content_layout = prs.slide_layouts[1]
        for slide_data in slides:
            s = prs.slides.add_slide(content_layout)
            s.shapes.title.text = _clean_text(slide_data.get("title", ""))
            if len(s.placeholders) > 1:
                tf = s.placeholders[1].text_frame
                tf.text = _clean_text(slide_data.get("content", ""))
                for para in tf.paragraphs:
                    for run in para.runs:
                        run.font.size = Pt(18)

        prs.save(str(out))
        logger.info(f"📊 PPTX generato: {out}")
        return {"success": True, "path": str(out), "filename": fname}

    except Exception as e:
        logger.error(f"Errore generazione PPTX: {e}")
        return {"error": str(e)}
